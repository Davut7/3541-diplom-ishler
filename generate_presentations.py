#!/usr/bin/env python3
"""Generate diploma project presentations for all 9 students."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
DARK_BG = RGBColor(15, 23, 42)
ACCENT = RGBColor(6, 182, 212)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(203, 213, 225)
MID_GRAY = RGBColor(148, 163, 184)
DARK_CARD = RGBColor(30, 41, 59)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(239, 68, 68)
YELLOW = RGBColor(250, 204, 21)
PURPLE = RGBColor(139, 92, 246)

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color=DARK_CARD, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, items, left, top, width, height, font_size=16, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def slide_title(slide, title, subtitle=None):
    set_slide_bg(slide)
    add_accent_line(slide, Inches(0.8), Inches(0.7), Inches(1.5))
    add_text(slide, title, Inches(0.8), Inches(0.8), Inches(8), Inches(0.8),
             font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, subtitle, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5),
                 font_size=16, color=MID_GRAY)

def create_presentation(project):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    name = project['student_name']
    proj_name = project['project_name']
    desc = project['one_line_description']
    problem = project['problem_statement']
    goals = project['goals']
    tech = project['tech_stack']
    features = project['key_features']
    how = project['how_it_works']
    arch = project['architecture']
    testing = project['testing_approach']

    # SLIDE 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(slide, proj_name, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             font_size=44, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1.5), Inches(3.2), Inches(3))
    add_text(slide, desc, Inches(1.5), Inches(3.5), Inches(9), Inches(0.8),
             font_size=20, color=MID_GRAY)
    add_text(slide, "Student: " + name, Inches(1.5), Inches(5.0), Inches(6), Inches(0.5),
             font_size=22, color=LIGHT_GRAY)
    add_text(slide, "Diploma Project  |  Cybersecurity  |  2026", Inches(1.5), Inches(5.5), Inches(6), Inches(0.5),
             font_size=16, color=MID_GRAY)

    # SLIDE 2: Table of Contents
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Table of Contents")
    toc = [
        "01  Introduction & Problem Statement",
        "02  Goals & Objectives",
        "03  Technology Stack",
        "04  System Architecture",
        "05  Key Features",
        "06  How It Works",
        "07  Implementation Details",
        "08  Testing & Results",
        "09  Future Work",
        "10  Conclusion"
    ]
    add_bullet_list(slide, toc, Inches(1.2), Inches(2.2), Inches(8), Inches(5), font_size=20)

    # SLIDE 3: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Introduction & Problem Statement")
    add_shape_bg(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(2.5))
    add_text(slide, "The Problem", Inches(1.2), Inches(2.5), Inches(4), Inches(0.5),
             font_size=20, color=ACCENT, bold=True)
    add_text(slide, problem, Inches(1.2), Inches(3.1), Inches(10.8), Inches(1.5),
             font_size=18, color=LIGHT_GRAY)
    add_shape_bg(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5))
    add_text(slide, "Why This Matters", Inches(1.2), Inches(5.4), Inches(4), Inches(0.5),
             font_size=20, color=YELLOW, bold=True)
    add_text(slide, "This project addresses a real-world cybersecurity challenge, providing practical tools and knowledge for defense against modern digital threats.",
             Inches(1.2), Inches(5.9), Inches(10.8), Inches(0.7), font_size=16, color=MID_GRAY)

    # SLIDE 4: Goals
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Goals & Objectives")
    for i, goal in enumerate(goals[:4]):
        y = Inches(2.3) + Inches(1.15) * i
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.95))
        num_color = [ACCENT, GREEN, PURPLE, YELLOW][i % 4]
        add_text(slide, "0" + str(i+1), Inches(1.1), y + Pt(8), Inches(0.6), Inches(0.5),
                 font_size=28, color=num_color, bold=True)
        add_text(slide, goal, Inches(1.8), y + Pt(12), Inches(10), Inches(0.7),
                 font_size=17, color=LIGHT_GRAY)

    # SLIDE 5: Tech Stack
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Technology Stack")
    cards_data = [
        ("Frontend", tech['frontend'], ACCENT),
        ("Backend", tech['backend'], GREEN),
        ("Key Libraries", tech['key_libraries'], PURPLE),
    ]
    card_w = Inches(3.5)
    gap = Inches(0.35)
    start_x = Inches(0.8)
    for i, (label, value, color) in enumerate(cards_data):
        x = start_x + (card_w + gap) * i
        add_shape_bg(slide, x, Inches(2.3), card_w, Inches(3.5))
        add_text(slide, label, x + Inches(0.3), Inches(2.5), Inches(3), Inches(0.5),
                 font_size=22, color=color, bold=True)
        add_accent_line(slide, x + Inches(0.3), Inches(3.1), Inches(1))
        add_text(slide, value, x + Inches(0.3), Inches(3.3), Inches(3), Inches(2),
                 font_size=16, color=LIGHT_GRAY)

    # SLIDE 6: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "System Architecture")
    add_text(slide, arch['description'], Inches(0.8), Inches(2.0), Inches(11.5), Inches(1),
             font_size=17, color=MID_GRAY)
    components = arch['components']
    cols = min(3, len(components))
    box_w = Inches(3.5)
    box_h = Inches(1.0)
    for i, comp in enumerate(components):
        col = i % cols
        row = i // cols
        x = Inches(0.8) + (box_w + Inches(0.35)) * col
        y = Inches(3.2) + (box_h + Inches(0.3)) * row
        add_shape_bg(slide, x, y, box_w, box_h)
        add_text(slide, comp, x + Inches(0.2), y + Inches(0.15), box_w - Inches(0.4), box_h - Inches(0.3),
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # SLIDE 7-8: Key Features
    half = len(features) // 2 + len(features) % 2
    for page, chunk in enumerate([features[:half], features[half:]]):
        if not chunk:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_text = "Key Features" if page == 0 else "Key Features (cont.)"
        slide_title(slide, title_text)
        for i, feat in enumerate(chunk):
            y = Inches(2.3) + Inches(0.95) * i
            add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.8))
            icon_colors = [ACCENT, GREEN, YELLOW, PURPLE, RED, ACCENT]
            add_text(slide, feat, Inches(1.5), y + Pt(8), Inches(10.5), Inches(0.6),
                     font_size=16, color=LIGHT_GRAY)

    # SLIDE 9: How It Works
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "How It Works")
    for i, step in enumerate(how):
        y = Inches(2.3) + Inches(1.15) * i
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.95))
        step_text = step.lstrip('0123456789. ')
        step_colors = [ACCENT, GREEN, YELLOW, PURPLE]
        add_text(slide, "Step " + str(i+1), Inches(1.1), y + Pt(6), Inches(1), Inches(0.5),
                 font_size=20, color=step_colors[i % 4], bold=True)
        add_text(slide, step_text, Inches(2.5), y + Pt(8), Inches(9.5), Inches(0.7),
                 font_size=16, color=LIGHT_GRAY)

    # SLIDE 10: Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Implementation Details")
    impl_items = [
        "Frontend: " + tech['frontend'] + " with responsive design and bilingual support (EN/TK)",
        "Backend: " + tech['backend'] + " with RESTful API architecture",
        "Dark/Light theme toggle with persistent user preferences",
        "Real-time data updates without page reload",
        "Key integrations: " + tech['key_libraries'],
    ]
    add_bullet_list(slide, impl_items, Inches(1.2), Inches(2.3), Inches(10.5), Inches(4.5), font_size=18)

    # SLIDE 11: Testing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Testing & Results")
    add_shape_bg(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(3))
    add_text(slide, "Testing Approach", Inches(1.2), Inches(2.5), Inches(4), Inches(0.5),
             font_size=22, color=GREEN, bold=True)
    add_text(slide, testing, Inches(1.2), Inches(3.2), Inches(10.8), Inches(1.8),
             font_size=17, color=LIGHT_GRAY)
    add_shape_bg(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.2))
    add_text(slide, "Results", Inches(1.2), Inches(5.9), Inches(4), Inches(0.5),
             font_size=22, color=ACCENT, bold=True)
    add_text(slide, "All security features demonstrated successfully with real data. The system correctly identifies threats and provides actionable security insights.",
             Inches(1.2), Inches(6.4), Inches(10.8), Inches(0.5), font_size=16, color=MID_GRAY)

    # SLIDE 12: Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Future Work")
    future = [
        "Integration with machine learning models for improved detection accuracy",
        "Support for additional file formats, protocols, or attack types",
        "Cloud deployment for multi-user access and scalability",
        "Mobile application version for on-the-go analysis",
        "Integration with third-party security tools and APIs",
    ]
    for i, item in enumerate(future):
        y = Inches(2.3) + Inches(0.95) * i
        add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(0.8))
        add_text(slide, item, Inches(1.6), y + Pt(8), Inches(10.5), Inches(0.6),
                 font_size=16, color=LIGHT_GRAY)

    # SLIDE 13: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title(slide, "Conclusion")
    add_shape_bg(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4))
    conclusions = [
        "Successfully developed " + proj_name + " addressing real cybersecurity challenges",
        "Implemented using modern web technologies: " + tech['frontend'] + " + " + tech['backend'],
        "Demonstrated practical security analysis with real data and interactive testing",
        "Provides educational value for understanding cybersecurity concepts",
        "Ready for further development and deployment in production environments",
    ]
    add_bullet_list(slide, conclusions, Inches(1.2), Inches(2.8), Inches(10.5), Inches(3), font_size=18)

    # SLIDE 14: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    add_text(slide, "Thank You!", Inches(1.5), Inches(2.0), Inches(10), Inches(1),
             font_size=52, color=WHITE, bold=True)
    add_accent_line(slide, Inches(1.5), Inches(3.2), Inches(3))
    add_text(slide, "Questions & Discussion", Inches(1.5), Inches(3.5), Inches(8), Inches(0.6),
             font_size=28, color=ACCENT)
    add_text(slide, "Presented by: " + name, Inches(1.5), Inches(4.8), Inches(6), Inches(0.5),
             font_size=20, color=LIGHT_GRAY)
    add_text(slide, "Project: " + proj_name, Inches(1.5), Inches(5.3), Inches(8), Inches(0.5),
             font_size=16, color=MID_GRAY)

    return prs


projects = [
    {
        "student_name": "Batyr Akyew",
        "folder": "batyr_akyew",
        "project_name": "WAF Behavioral Analysis",
        "one_line_description": "A Web Application Firewall combining signature-based detection with behavioral analysis to protect web applications.",
        "problem_statement": "Web applications face constant threats from SQL injection, XSS, path traversal, and command injection. Traditional WAFs use only static rules and miss sophisticated attack patterns.",
        "goals": ["Detect and block SQL Injection, XSS, Path Traversal, and Command Injection attacks", "Implement behavioral analysis to identify bot traffic and anomalous sessions", "Provide real-time monitoring dashboard with threat scoring and attack logging", "Demonstrate how signature-based and behavioral detection complement each other"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js, SQLite", "key_libraries": "sql.js, JWT, bcrypt, PDFKit, Rate Limiter"},
        "key_features": ["SQL Injection, XSS, Path Traversal, Command Injection detection", "User session monitoring and bot detection via behavioral analysis", "Real-time threat scoring with 4 risk levels", "Live attack logs and statistics dashboard with charts", "Interactive request analyzer for testing payloads", "PDF security report generation"],
        "how_it_works": ["Incoming HTTP requests are intercepted by the WAF middleware", "Each request is checked against signature-based rules (regex patterns)", "Behavioral analysis tracks session patterns, request frequency, and anomalies", "Based on combined risk score (0-100%), requests are allowed, monitored, or blocked"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "WAF Rules Engine", "Behavioral Analysis Module", "SQLite Database", "PDF Report Generator"], "description": "Two-tier architecture: Vue.js frontend communicates via REST API with Express.js backend containing the WAF engine, behavioral profiler, and logging system."},
        "testing_approach": "Real attack payloads are sent through the WAF engine on startup. The Analyzer page allows interactive testing. All statistics come from genuine WAF analysis results."
    },
    {
        "student_name": "Bayram Gochmyradow",
        "folder": "bayram_gochmyradow",
        "project_name": "Android Security Analyzer",
        "one_line_description": "A web-based tool for analyzing Android APK file security by examining permissions, manifest vulnerabilities, and risk scores.",
        "problem_statement": "Android users and developers lack a simple tool to analyze APK files for dangerous permissions, unsafe manifest settings, and potential security risks before installation.",
        "goals": ["Analyze Android APK files for dangerous permissions across 70+ types", "Check AndroidManifest.xml for insecure configurations", "Calculate comprehensive risk scores with actionable recommendations", "Provide exportable reports in JSON and HTML formats"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js, Multer", "key_libraries": "adm-zip, xml2js, uuid"},
        "key_features": ["Drag-and-drop APK upload with real-time analysis", "Permission analysis with 4 risk levels: Critical, High, Medium, Low", "AndroidManifest.xml vulnerability scanning", "Risk score calculation (0-100%) with weighted formula", "Report export in JSON and HTML formats", "Analysis history with statistics charts"],
        "how_it_works": ["User uploads APK file via drag-and-drop interface", "Backend extracts AndroidManifest.xml from APK using adm-zip", "Manifest is parsed and permissions checked against 70+ known types", "Detailed report generated with risk score, issues, and recommendations"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "APK Analyzer Core", "Permission Database (70+)", "Multer File Handler", "Report Generator"], "description": "Two-tier architecture: Frontend handles file upload and results. Backend unpacks APK, parses manifest, runs analysis, and returns structured results."},
        "testing_approach": "Real APK files can be uploaded for live analysis. The system extracts actual AndroidManifest.xml and checks against a comprehensive permission database."
    },
    {
        "student_name": "Daniyar Nurmedow",
        "folder": "daniyar_nurmedow",
        "project_name": "XSS Shield - Attack & Defense Laboratory",
        "one_line_description": "An educational platform demonstrating XSS attack techniques and defense methods with interactive code scanning.",
        "problem_statement": "Cross-Site Scripting (XSS) remains one of the most widespread web vulnerabilities (OWASP Top 10). Developers need a safe environment to understand and defend against all XSS types.",
        "goals": ["Demonstrate Reflected, Stored, and DOM-based XSS attacks safely", "Teach defense techniques: input validation, output encoding, CSP, HttpOnly cookies", "Provide a code scanner detecting XSS vulnerability patterns", "Show real-world XSS payload examples and countermeasures"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js", "key_libraries": "DOMPurify, XSS Pattern Scanner, uuid"},
        "key_features": ["Interactive XSS attack simulator for 3 XSS types", "Source code scanner detecting dangerous patterns (innerHTML, eval, v-html)", "Defense technique demos with live code examples", "Real-world XSS payload library", "CSP configuration examples", "Persistent scan history"],
        "how_it_works": ["User selects an attack type or enters code to scan", "Attack simulator demonstrates XSS payload execution in controlled environment", "Code scanner checks against known vulnerability patterns and rates severity", "Defense recommendations provided with working code examples"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "XSS Attack Simulator", "Code Vulnerability Scanner", "Scan History (JSON)", "Defense Examples Engine"], "description": "Two-tier: Frontend provides interactive attack/defense lab. Backend handles code scanning, manages history, and serves simulation data."},
        "testing_approach": "Users test real XSS payloads in the simulator and scan source code. Scanner detects patterns across frameworks (jQuery, Vue, React). History persisted for review."
    },
    {
        "student_name": "Dawutmuhammet Begmedow",
        "folder": "dawutmuhammet_begmedow",
        "project_name": "VirusDetect Pro - Advanced Malware Detection",
        "one_line_description": "A malware detection system using heuristic analysis, behavioral detection, and entropy analysis to identify viruses beyond signatures.",
        "problem_statement": "Traditional antivirus relies on known signatures and fails to detect polymorphic, metamorphic, packed, fileless, or rootkit malware that evades signature-based detection.",
        "goals": ["Demonstrate detection beyond signatures: heuristic, behavioral, entropy analysis", "Explain evasion techniques used by modern malware", "Provide file upload with SHA-256 hashing and multi-method scanning", "Educate on why traditional antivirus is insufficient"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js, Multer, crypto", "key_libraries": "crypto (SHA-256), better-sqlite3, adm-zip, Multer"},
        "key_features": ["File upload with SHA-256 hash computation", "Heuristic analysis for suspicious API calls", "Behavioral detection in sandbox simulation", "Entropy analysis for packed/encrypted malware", "Malware signature database", "Persistent scan history with SQLite"],
        "how_it_works": ["User uploads a suspicious file", "System computes SHA-256 hash and checks against signatures", "Heuristic and behavioral analysis scan for suspicious patterns", "Entropy analysis checks for packing; results compiled into threat report"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "Signature Database", "Heuristic Analysis Engine", "Behavioral Detection Module", "Entropy Calculator", "SQLite Database"], "description": "Two-tier with SQLite persistence: Frontend provides upload UI. Backend performs multi-method file analysis and stores results."},
        "testing_approach": "Files uploaded for real-time analysis. System computes actual SHA-256 hashes and runs pattern matching. Scan history persisted in SQLite."
    },
    {
        "student_name": "Rowshen Palwanow",
        "folder": "rowshen_palwanow",
        "project_name": "KeyGuard - Keylogger Detection",
        "one_line_description": "A keylogger detection tool monitoring processes, hooks, and API calls to identify hidden keylogging threats.",
        "problem_statement": "Keyloggers are among the most dangerous hidden threats, stealing passwords and personal data. Users need specialized tools to detect software keyloggers.",
        "goals": ["Detect software keyloggers using hook monitoring (SetWindowsHookEx)", "Provide real-time process analysis for suspicious behavior", "Monitor registry for persistence mechanisms", "Educate on keylogger types and protection strategies"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js, child_process", "key_libraries": "child_process (execSync), os module, uuid"},
        "key_features": ["Hook detection for SetWindowsHookEx API calls", "API monitoring for GetAsyncKeyState and keybd_event", "Process analysis identifying suspicious behaviors", "Registry scanning for persistence mechanisms", "Network monitoring for data exfiltration", "Safe process whitelist"],
        "how_it_works": ["System scans running processes against whitelist", "Hook detection checks for keyboard hook APIs", "API and registry monitoring identifies suspicious patterns", "Results displayed with risk assessment and recommendations"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "Process Scanner", "Hook Detection Module", "Registry Scanner", "Network Monitor", "Scan History (JSON)"], "description": "Two-tier: Frontend displays detection dashboard. Backend uses child_process and os modules to scan processes and detect keyloggers."},
        "testing_approach": "Real process scanning using OS commands. Whitelist reduces false positives. Scan history persisted. System runs actual process checks."
    },
    {
        "student_name": "Selbi Weliyyewa",
        "folder": "selbi_weliyyewa",
        "project_name": "GAN Security Analyzer",
        "one_line_description": "A platform analyzing 9 adversarial attack types and 7 defense mechanisms against Generative Adversarial Networks.",
        "problem_statement": "GANs are deployed in critical applications but their security vulnerabilities are poorly understood. Understanding adversarial attacks and defenses is essential for safe AI.",
        "goals": ["Document 9 types of adversarial attacks against GANs", "Present 7 defense mechanisms with effectiveness metrics", "Provide interactive attack simulator with adjustable parameters", "Compile research statistics from 2014-2026 on AI security"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js", "key_libraries": "Chart.js, PrimeVue 4"},
        "key_features": ["9 attack types with formulas and research data", "7 defense mechanisms with effectiveness ratings", "Interactive attack simulator", "Statistical analysis with 4 chart types", "Attack vs defense comparison tables", "Bilingual interface (EN/TK)"],
        "how_it_works": ["User explores 9 attack types with descriptions and formulas", "Simulator allows selecting attack type and adjusting strength", "System computes perturbation, confidence drop, success probability", "Statistics visualize trends and defense effectiveness"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "Attack Database (9 types)", "Defense Database (7 mechanisms)", "Attack Simulator Engine", "Statistics Module"], "description": "Two-tier: Frontend provides interactive exploration. Backend serves research data, runs simulations, and provides analytics."},
        "testing_approach": "Interactive simulator with adjustable parameters. Results include perturbation values and defense recommendations. Based on real published research."
    },
    {
        "student_name": "Shanur Gulmyradow",
        "folder": "shanur_gulmyradow",
        "project_name": "Wireshark Network Monitor",
        "one_line_description": "A web-based network traffic analysis tool providing packet capture, protocol analysis, and security monitoring.",
        "problem_statement": "Network administrators need accessible tools to analyze traffic, detect threats, and understand protocol behavior without deep CLI expertise.",
        "goals": ["Simulate Wireshark-like packet capture in a web interface", "Support 12+ network protocols across OSI layers", "Provide security monitoring with anomaly detection", "Visualize traffic patterns with real-time charts"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js", "key_libraries": "Chart.js, PrimeVue 4"},
        "key_features": ["Real-time packet capture with Wireshark-like interface", "Protocol analysis for 12+ protocols (HTTP, DNS, SSH, etc.)", "OSI model visualization and protocol hierarchy", "Security monitoring: port scan, DDoS, exfiltration detection", "Traffic statistics with distribution charts", "BPF-like capture filters"],
        "how_it_works": ["User selects network interface and applies filters", "System captures and generates realistic network packets", "Packets decoded through protocol dissectors", "Security monitoring analyzes patterns and generates alerts"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "Packet Capture Engine", "Protocol Dissector", "Traffic Analysis Engine", "Security Alert System", "Statistics Module"], "description": "Two-tier: Frontend displays packet list, protocol details, hex dump. Backend generates packets, analyzes protocols, runs anomaly detection."},
        "testing_approach": "Packet capture generates realistic traffic. Users apply filters, inspect packets, view protocols. Statistics computed in real-time."
    },
    {
        "student_name": "Shatlyk Rahmanov",
        "folder": "shatlyk_rahmanov",
        "project_name": "AI Firewall - Intelligent Network Protection",
        "one_line_description": "An AI-powered firewall using neural networks for threat detection and automatic firewall rule generation.",
        "problem_statement": "Traditional firewalls rely on static rules and cannot adapt to new threats. They fail to detect zero-day attacks and require constant manual management.",
        "goals": ["Automate firewall rule generation using neural network analysis", "Detect zero-day attacks through behavioral analysis", "Reduce false positives vs traditional firewalls", "Provide real-time AI-powered threat classification"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js", "backend": "Node.js, Express.js", "key_libraries": "Chart.js, Neural Network simulation engine"},
        "key_features": ["Neural network: 6-input, two hidden layers (128/64), 3-class output", "Automatic firewall rule generation", "Real-time traffic classification (Safe/Suspicious/Threat)", "Zero-day threat detection via anomaly analysis", "DDoS, port scan, brute force, malware C&C detection", "AI model status dashboard"],
        "how_it_works": ["Traffic features extracted (packet size, protocol, ports, duration)", "Neural network processes through hidden layers with ReLU", "Traffic classified as Safe, Suspicious, or Threat", "AI suggests or generates firewall rules based on patterns"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "Neural Network Engine", "Traffic Monitor", "Threat Detection Module", "Auto Rule Generator", "AI Learning Module"], "description": "Two-tier: Frontend for rule management, monitoring, AI status. Backend runs neural network classification and generates rules."},
        "testing_approach": "AI endpoint accepts traffic data for classification. Demonstrates rule suggestion and generation. AI model status viewable in dashboard."
    },
    {
        "student_name": "Suleyman Akmuhammedow",
        "folder": "suleyman_akmuhammedow",
        "project_name": "OSINT.AI - Open Source Intelligence Tools",
        "one_line_description": "A web-based OSINT platform for DNS reconnaissance, IP geolocation, port scanning, WHOIS lookup, and AI risk assessment.",
        "problem_statement": "Cybersecurity professionals need a unified tool to gather and analyze publicly available intelligence to assess security posture proactively and ethically.",
        "goals": ["Provide comprehensive OSINT: DNS, WHOIS, port scanning, geolocation", "Calculate AI-powered risk scores from discovered services", "Generate exportable reports in JSON and HTML", "Ensure ethical, non-destructive, read-only operations"],
        "tech_stack": {"frontend": "Vue.js 3, PrimeVue 4, Chart.js, Axios", "backend": "Node.js, Express.js, DNS module", "key_libraries": "dns (Node.js built-in), child_process, Chart.js"},
        "key_features": ["DNS reconnaissance and record discovery", "IP geolocation with ISP identification", "Port scanning and service identification", "WHOIS domain registration data", "AI-powered risk score (weighted by port criticality)", "Report export in JSON and HTML"],
        "how_it_works": ["User enters target domain or IP address", "System performs DNS, ping, port scan, WHOIS, and GeoIP lookup", "AI calculates risk score (Low/Medium/High/Critical)", "Detailed report generated with findings and export options"],
        "architecture": {"components": ["Vue.js 3 SPA Frontend", "Express.js REST API", "DNS Resolution Module", "Port Scanner", "WHOIS Module", "GeoIP Module", "Risk Assessment Algorithm", "Report Generator"], "description": "Two-tier: Frontend for target input and results. Backend performs network reconnaissance using DNS module and child_process."},
        "testing_approach": "Real domains/IPs can be analyzed live. All operations read-only and non-destructive. Report history with export capabilities."
    }
]

output_dir = '/Users/macbookpro/Desktop/projects/diplom_ishler/3541/presentations'
os.makedirs(output_dir, exist_ok=True)

for project in projects:
    prs = create_presentation(project)
    filename = project['folder'] + "_presentation.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    print("Done: " + filename)

print("\nAll 9 presentations saved to " + output_dir)
